from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import pandas as pd


def extract_info_from_excel(input_excel):
    try:
        # Read the Excel file without headers
        #df = pd.read_excel(input_excel, sheet_name='MainConference', header=None)
        #df = pd.read_excel(input_excel, sheet_name='MainConference_Workshop', header=None)
        #df = pd.read_excel(input_excel, sheet_name='MainConference_Workshop', header=None)
        #df = pd.read_excel(input_excel, sheet_name='Guests', header=None)
        df = pd.read_excel(input_excel, sheet_name='Organizers', header=None)
        #df = pd.read_excel(input_excel, sheet_name='Workshops', header=None)

        # Extract required information from the columns
        info_list = []
        for index, row in df.iterrows():
            # Extract data from the columns by their index
            first_name = row[0]
            last_name = row[1]
            affiliation = row[3]
            country = row[2]

            # Combine first and last name
            full_name = f"{first_name} {last_name}"

            # Append extracted information to the list
            info_list.append((full_name, affiliation, country))

        return info_list
    except Exception as e:
        print("An error occurred:", e)
        return []


def find_replace_text(pptx_file, replace_list, output_file):
    replace_full_name = replace_list[0]
    replace_affiliation = replace_list[1]
    replace_country = replace_list[2]
    # Load the PowerPoint presentation
    prs = Presentation(pptx_file)

    # Iterate through each slide in the presentation
    for slide in prs.slides:
        # Iterate through each shape (text box) in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Check if the shape contains the search text
                if 'NAME' in shape.text:
                    # Copy font properties
                    font = shape.text_frame.paragraphs[0].runs[0].font
                    # Replace text
                    shape.text = shape.text.replace('NAME', replace_full_name)
                    # Apply font properties to replacement text
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            #run.font.font_name = font.font_name
                            if 'Kariyakarawana' in replace_full_name:
                                run.font.size = font.size - 3
                            else:
                                run.font.size = font.size
                            run.font.bold = font.bold

                elif 'AFFILIATION' in shape.text:
                    # Copy font properties
                    font = shape.text_frame.paragraphs[0].runs[0].font
                    # Replace text
                    shape.text = shape.text.replace('AFFILIATION', replace_affiliation)
                    # Apply font properties to replacement text
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            #run.font.font_name =  font.font_name
                            run.font.size = font.size
                            run.font.bold = font.bold


                elif 'COUNTRY' in shape.text:
                    # Copy font properties
                    font = shape.text_frame.paragraphs[0].runs[0].font
                    # Replace text
                    shape.text = shape.text.replace('COUNTRY', replace_country)
                    # Apply font properties to replacement text
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            #run.font.font_name = font.font_name
                            run.font.size = font.size
                            run.font.bold = font.bold

    # Save the modified presentation to the output file
    prs.save(output_file)

#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\conference_badge_ppt.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\conference_badge_ppt_empty.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\workshop.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\conference_workshop_badge.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\organizer_badge.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\guests_badge_empty.pptx'  #'badge_ppt.pptx'
pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\conference_badge_ppt_empty.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\conference_workshop_badge_empty.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\workshop_empty.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\organizer_badge_empty.pptx'  #'badge_ppt.pptx'
#pptx_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\guests_badge.pptx'  #'badge_ppt.pptx'
search_text = 'NAME'
replace_text = 'FIRST_NAME LAST_NAME'
input_excel = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\SANER2024_Registered.xlsx'
output_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\intermediate_conference_badges\\'

#output_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\intermediate_workshop_badges\\'
#output_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\intermediate_conf_workshop_badges\\'
#output_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\intermediate_organizer_badges\\'
#output_file = 'C:\\Users\\msridhar20\\PycharmProjects\\scopus_data_processing\\badges\\intermediate_guest_badges\\'
info = extract_info_from_excel(input_excel)
cntr = 0
for item in info:
    cntr += 1
    print("Name:", item[0])
    print("Affiliation:", item[1])
    print("Country:", item[2])
    print()
    find_replace_text(pptx_file, item, output_file+item[0]+'_'+item[2]+'.pptx')
    #find_replace_text(pptx_file, item, output_file + 'empty_workshop_backup_badge_' + str(cntr) + '.pptx')
